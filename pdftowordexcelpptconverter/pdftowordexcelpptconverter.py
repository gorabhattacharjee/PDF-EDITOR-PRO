import os
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook
from PIL import Image
import pytesseract

# ---- Paths - update if necessary ----
POPPLER_PATH = r"C:\poppler\Library\bin"
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# ---- Converters ----

def pdf_to_word(pdf_path, output_docx="output.docx"):
    document = Document()
    pages = convert_from_path(pdf_path, 300, poppler_path=POPPLER_PATH)

    for i, page in enumerate(pages, start=1):
        img_path = f"temp_page_{i}.png"
        page.save(img_path, "PNG")

        document.add_picture(img_path, width=Inches(6.5))
        text = pytesseract.image_to_string(page)
        document.add_paragraph(text)

        os.remove(img_path)
        if i < len(pages):
            document.add_page_break()

    document.save(output_docx)
    print(f"✅ Word exported → {output_docx}")


def pdf_to_excel(pdf_path, output_xlsx="output.xlsx"):
    wb = Workbook()
    ws = wb.active
    ws.title = "PDF to Excel"

    pages = convert_from_path(pdf_path, 200, poppler_path=POPPLER_PATH)

    for i, page in enumerate(pages, start=1):
        text = pytesseract.image_to_string(page)
        for row in text.split("\n"):
            ws.append([row])

    wb.save(output_xlsx)
    print(f"✅ Excel exported → {output_xlsx}")


def pdf_to_ppt(pdf_path, output_ppt="output.pptx"):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    pages = convert_from_path(pdf_path, 300, poppler_path=POPPLER_PATH)

    for i, page in enumerate(pages, start=1):
        slide = prs.slides.add_slide(blank_layout)
        img_path = f"temp_page_{i}.png"
        page.save(img_path, "PNG")

        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width)
        os.remove(img_path)

    prs.save(output_ppt)
    print(f"✅ PowerPoint exported → {output_ppt}")


# ---- Main Execution ----
if __name__ == "__main__":
    pdf_file = "sample.pdf"  # Change to your file path

    pdf_to_word(pdf_file, "converted.docx")
    pdf_to_excel(pdf_file, "converted.xlsx")
    pdf_to_ppt(pdf_file, "converted.pptx")

    print("✅✅✅ All Conversions Completed Successfully!")
