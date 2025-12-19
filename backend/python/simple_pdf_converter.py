#!/usr/bin/env python3
"""
Simple, fast PDF to Office conversion script
Outputs binary data directly to stdout for backend streaming
"""
import sys
import os

def pdf_to_word_simple(pdf_path, output_docx):
    """Simple PDF to Word conversion using pdf2docx with proper page sizing"""
    try:
        from pdf2docx import Converter
        from docx import Document
        from docx.shared import Inches, Pt
        import fitz
        
        print(f"[pdf_to_word] Starting conversion: {pdf_path}", file=sys.stderr)
        print(f"[pdf_to_word] Output: {output_docx}", file=sys.stderr)
        
        # Step 1: Get PDF dimensions
        print(f"[pdf_to_word] Measuring PDF page size...", file=sys.stderr)
        pdf_doc = fitz.open(pdf_path)
        if len(pdf_doc) > 0:
            page = pdf_doc[0]
            rect = page.rect
            # Convert from points to inches (72 points = 1 inch)
            pdf_width_inches = rect.width / 72.0
            pdf_height_inches = rect.height / 72.0
            print(f"[pdf_to_word] PDF page size: {pdf_width_inches:.2f}\" x {pdf_height_inches:.2f}\"", file=sys.stderr)
        else:
            pdf_width_inches = 8.5  # Default US Letter
            pdf_height_inches = 11.0
            print(f"[pdf_to_word] Using default page size: {pdf_width_inches:.2f}\" x {pdf_height_inches:.2f}\"", file=sys.stderr)
        pdf_doc.close()
        
        # Step 2: Convert PDF to Word using pdf2docx
        print(f"[pdf_to_word] Running conversion...", file=sys.stderr)
        cv = Converter(pdf_path)
        cv.convert(output_docx, multi_processing=False, cpu_count=1)
        cv.close()
        
        # Step 3: Adjust page size in Word document to match PDF
        print(f"[pdf_to_word] Adjusting page size in Word...", file=sys.stderr)
        doc = Document(output_docx)
        
        # Minimum margins in inches
        min_margin = 0.25
        
        # Set page size for all sections
        for section in doc.sections:
            # Set page width and height to PDF size
            section.page_width = Inches(pdf_width_inches)
            section.page_height = Inches(pdf_height_inches)
            
            # Set margins to minimum
            section.top_margin = Inches(min_margin)
            section.bottom_margin = Inches(min_margin)
            section.left_margin = Inches(min_margin)
            section.right_margin = Inches(min_margin)
            
            print(f"[pdf_to_word] Set page size: {pdf_width_inches:.2f}\" x {pdf_height_inches:.2f}\" with {min_margin}\" margins", file=sys.stderr)
        
        # Save the modified document
        doc.save(output_docx)
        
        # Verify output
        if os.path.exists(output_docx) and os.path.getsize(output_docx) > 0:
            size = os.path.getsize(output_docx)
            print(f"[pdf_to_word] SUCCESS: Created {size} bytes", file=sys.stderr)
            return True
        else:
            print(f"[pdf_to_word] ERROR: Output file is empty or missing", file=sys.stderr)
            return False
    except Exception as e:
        print(f"[pdf_to_word] EXCEPTION: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc(file=sys.stderr)
        return False

def pdf_to_excel_simple(pdf_path, output_xlsx):
    """Simple PDF to Excel conversion - extracts ALL pages and content"""
    try:
        import pdfplumber
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        
        print(f"[pdf_to_excel] Starting conversion: {pdf_path}", file=sys.stderr)
        print(f"[pdf_to_excel] Output: {output_xlsx}", file=sys.stderr)
        
        wb = Workbook()
        wb.remove(wb.active)  # Remove default sheet
        
        total_sheets_created = 0
        total_tables_found = 0
        
        with pdfplumber.open(pdf_path) as pdf:
            total_pages = len(pdf.pages)
            print(f"[pdf_to_excel] PDF has {total_pages} pages total", file=sys.stderr)
            
            for page_idx in range(total_pages):
                page = pdf.pages[page_idx]
                page_number = page_idx + 1
                
                print(f"[pdf_to_excel] ===== PROCESSING PAGE {page_number}/{total_pages} =====", file=sys.stderr)
                
                # Try to extract tables
                try:
                    tables = page.extract_tables()
                    print(f"[pdf_to_excel] page.extract_tables() returned: {len(tables) if tables else 0} tables", file=sys.stderr)
                except Exception as e:
                    print(f"[pdf_to_excel] Error extracting tables from page {page_number}: {e}", file=sys.stderr)
                    tables = None
                
                # If tables found, create sheets for them
                if tables and len(tables) > 0:
                    for table_idx, table in enumerate(tables, 1):
                        sheet_name = f"P{page_number}_T{table_idx}"
                        ws = wb.create_sheet(title=sheet_name)
                        
                        print(f"[pdf_to_excel] Creating sheet '{sheet_name}' with {len(table)} rows", file=sys.stderr)
                        
                        # Write all rows from table
                        for row_idx, row in enumerate(table, 1):
                            for col_idx, cell_value in enumerate(row, 1):
                                value = cell_value if cell_value else ""
                                ws.cell(row=row_idx, column=col_idx, value=value)
                        
                        total_tables_found += 1
                        total_sheets_created += 1
                
                # Always extract text content from page (even if tables exist)
                try:
                    text = page.extract_text()
                    if text and text.strip():
                        print(f"[pdf_to_excel] Extracted text from page {page_number} ({len(text)} chars)", file=sys.stderr)
                        
                        # Create text sheet
                        text_sheet_name = f"P{page_number}_Text"
                        ws_text = wb.create_sheet(title=text_sheet_name)
                        
                        # Split text into lines
                        lines = text.split('\n')
                        print(f"[pdf_to_excel] Adding {len(lines)} lines to text sheet", file=sys.stderr)
                        
                        for row_idx, line in enumerate(lines, 1):
                            ws_text.cell(row=row_idx, column=1, value=line)
                        
                        # Set column width
                        ws_text.column_dimensions['A'].width = 100
                        total_sheets_created += 1
                    else:
                        print(f"[pdf_to_excel] No text content on page {page_number}", file=sys.stderr)
                except Exception as e:
                    print(f"[pdf_to_excel] Error extracting text from page {page_number}: {e}", file=sys.stderr)
        
        # Save the workbook
        print(f"[pdf_to_excel] Creating {len(wb.sheetnames)} sheets...", file=sys.stderr)
        print(f"[pdf_to_excel] Sheet names: {wb.sheetnames}", file=sys.stderr)
        wb.save(output_xlsx)
        
        # Verify output
        if os.path.exists(output_xlsx) and os.path.getsize(output_xlsx) > 0:
            size = os.path.getsize(output_xlsx)
            print(f"[pdf_to_excel] ✅ SUCCESS: {size} bytes", file=sys.stderr)
            print(f"[pdf_to_excel] ✅ Created {len(wb.sheetnames)} sheets from {total_pages} pages", file=sys.stderr)
            print(f"[pdf_to_excel] ✅ Found {total_tables_found} tables total", file=sys.stderr)
            return True
        else:
            print(f"[pdf_to_excel] ❌ ERROR: Output file empty or missing", file=sys.stderr)
            return False
    except Exception as e:
        print(f"[pdf_to_excel] ❌ EXCEPTION: {str(e)}", file=sys.stderr)
        import traceback
        traceback.print_exc(file=sys.stderr)
        return False

def pdf_to_ppt_simple(pdf_path, output_pptx):
    """Simple PDF to PowerPoint conversion"""
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches
        import io
        
        print(f"[pdf_to_ppt] Starting conversion: {pdf_path}", file=sys.stderr)
        print(f"[pdf_to_ppt] Output: {output_pptx}", file=sys.stderr)
        
        # Convert PDF pages to images
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Try with poppler path if on Windows
        poppler_path = None
        if os.name == 'nt':  # Windows
            try:
                import glob
                user_local = os.path.join(os.environ.get('LOCALAPPDATA', ''), 'Microsoft', 'WinGet', 'Packages')
                poppler_dirs = glob.glob(os.path.join(user_local, 'oschwartz10612.Poppler*', 'poppler-*', 'Library', 'bin'))
                if poppler_dirs:
                    poppler_path = poppler_dirs[0]
                    print(f"[pdf_to_ppt] Using poppler at: {poppler_path}", file=sys.stderr)
            except:
                pass
        
        pages = convert_from_path(pdf_path, poppler_path=poppler_path)
        
        print(f"[pdf_to_ppt] Converting {len(pages)} pages to PowerPoint...", file=sys.stderr)
        
        for idx, page_image in enumerate(pages, 1):
            # Save image to bytes
            img_byte_arr = io.BytesIO()
            page_image.save(img_byte_arr, format='PNG')
            img_byte_arr.seek(0)
            
            # Add slide with image
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            left = Inches(0)
            top = Inches(0)
            pic = slide.shapes.add_picture(img_byte_arr, left, top, width=prs.slide_width, height=prs.slide_height)
            
            print(f"[pdf_to_ppt] Added page {idx}/{len(pages)}", file=sys.stderr)
        
        prs.save(output_pptx)
        
        if os.path.exists(output_pptx) and os.path.getsize(output_pptx) > 0:
            size = os.path.getsize(output_pptx)
            print(f"[pdf_to_ppt] SUCCESS: Created {size} bytes", file=sys.stderr)
            return True
        else:
            print(f"[pdf_to_ppt] ERROR: Output file is empty or missing", file=sys.stderr)
            return False
    except Exception as e:
        print(f"[pdf_to_ppt] EXCEPTION: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc(file=sys.stderr)
        return False

def pdf_to_html_simple(pdf_path, output_html):
    """Simple PDF to HTML conversion"""
    try:
        import pdfplumber
        
        print(f"[pdf_to_html] Starting conversion: {pdf_path}", file=sys.stderr)
        print(f"[pdf_to_html] Output: {output_html}", file=sys.stderr)
        
        html_content = '<html><head><meta charset="utf-8"><title>PDF to HTML</title></head><body>'
        
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                html_content += f'<h2>Page {page_num}</h2>'
                
                # Extract text
                text = page.extract_text()
                if text:
                    html_content += f'<p>{text.replace(chr(10), "<br>")}</p>'
                
                # Extract tables
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        html_content += '<table border="1">'
                        for row in table:
                            html_content += '<tr>'
                            for cell in row:
                                html_content += f'<td>{cell or ""}</td>'
                            html_content += '</tr>'
                        html_content += '</table>'
        
        html_content += '</body></html>'
        
        with open(output_html, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        if os.path.exists(output_html) and os.path.getsize(output_html) > 0:
            size = os.path.getsize(output_html)
            print(f"[pdf_to_html] SUCCESS: Created {size} bytes", file=sys.stderr)
            return True
        else:
            print(f"[pdf_to_html] ERROR: Output file is empty or missing", file=sys.stderr)
            return False
    except Exception as e:
        print(f"[pdf_to_html] EXCEPTION: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc(file=sys.stderr)
        return False

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python simple_pdf_converter.py <format> <input_pdf> <output_file>", file=sys.stderr)
        print("Formats: word, excel, ppt, html", file=sys.stderr)
        sys.exit(1)
    
    format_type = sys.argv[1].lower()
    input_pdf = sys.argv[2]
    output_file = sys.argv[3]
    
    print(f"[Main] Format: {format_type}", file=sys.stderr)
    print(f"[Main] Input: {input_pdf}", file=sys.stderr)
    print(f"[Main] Output: {output_file}", file=sys.stderr)
    
    try:
        success = False
        
        if format_type == "word":
            success = pdf_to_word_simple(input_pdf, output_file)
        elif format_type == "excel":
            success = pdf_to_excel_simple(input_pdf, output_file)
        elif format_type == "ppt":
            success = pdf_to_ppt_simple(input_pdf, output_file)
        elif format_type == "html":
            success = pdf_to_html_simple(input_pdf, output_file)
        else:
            print(f"[Main] Unknown format: {format_type}", file=sys.stderr)
            sys.exit(1)
        
        if success:
            # Read and output file to stdout
            print(f"[Main] Reading output file...", file=sys.stderr)
            with open(output_file, 'rb') as f:
                output_data = f.read()
            
            print(f"[Main] Sending {len(output_data)} bytes to stdout", file=sys.stderr)
            sys.stdout.buffer.write(output_data)
            print(f"[Main] DONE", file=sys.stderr)
        else:
            print(f"[Main] Conversion failed", file=sys.stderr)
            sys.exit(1)
    
    except Exception as e:
        print(f"[Main] EXCEPTION: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)
